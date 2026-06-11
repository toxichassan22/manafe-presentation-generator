"""
PPTX Builder – redesigned to match reference design exactly.
- Cover / Section headers: full-bleed image + dark overlay + triangle decorations + centered text
- Standard slides: white background + burgundy title + gray triangle decorations + 3-part footer
"""

import io
import logging
from datetime import datetime
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from config.settings import settings, PROJECT_ROOT
from config.branding import branding

logger = logging.getLogger(__name__)

# ── Constants ─────────────────────────────────────────────────────────
TRIANGLE = 7          # msoShapeIsoscelesTriangle
RECTANGLE = 1         # msoShapeRectangle


def _get_color(hex_str: str) -> RGBColor:
    hex_str = hex_str.lstrip('#')
    return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))


def _mark_rtl(paragraph):
    """Mark a paragraph as RTL so Arabic text renders correctly."""
    try:
        paragraph._p.get_or_add_pPr().set("rtl", "1")
    except Exception:
        pass


def set_shape_transparency(shape, alpha_val):
    """Set the transparency/alpha value (0-100) of a shape's fill using XML manipulation."""
    try:
        fill = shape.fill
        if alpha_val == 0:
            return
        spPr = shape._element.spPr
        solidFill = spPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
        if solidFill is None:
            fill.solid()
            solidFill = spPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')

        if solidFill is not None:
            color_elem = None
            for tag in ['srgbClr', 'schemeClr']:
                color_elem = solidFill.find(f'{{http://schemas.openxmlformats.org/drawingml/2006/main}}{tag}')
                if color_elem is not None:
                    break

            if color_elem is not None:
                alpha_elem = color_elem.find('{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
                if alpha_elem is None:
                    from lxml import etree
                    alpha_elem = etree.SubElement(
                        color_elem,
                        '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha'
                    )
                alpha_elem.set('val', str(int((100 - alpha_val) * 1000)))
    except Exception as e:
        logger.warning("Could not set shape transparency: %s", e)


def _set_slide_bg(slide, color: RGBColor):
    """Set solid background color for a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_image(slide, prs, image_bytes: bytes, left: float = 0.0, top: float = 0.0,
               width: float = None, height: float = None):
    """Add an image from bytes to the slide."""
    stream = io.BytesIO(image_bytes)
    w = Inches(width) if width is not None else prs.slide_width
    h = Inches(height) if height is not None else prs.slide_height
    slide.shapes.add_picture(stream, Inches(left), Inches(top), w, h)


def _add_image_path(slide, prs, image_path: Path, left: float = 0.0, top: float = 0.0,
                    width: float = None, height: float = None):
    """Add an image from a file path to the slide."""
    w = Inches(width) if width is not None else prs.slide_width
    h = Inches(height) if height is not None else prs.slide_height
    slide.shapes.add_picture(str(image_path), Inches(left), Inches(top), w, h)



# ── Decorative Triangles ──────────────────────────────────────────────

def _add_triangles_cover(slide, prs):
    """Add triangle decorations matching the reference cover/section style - all within bounds."""
    dark_red  = RGBColor(55, 8, 8)
    mid_red   = RGBColor(80, 15, 15)
    white     = RGBColor(255, 255, 255)
    W = prs.slide_width   # in EMUs
    H = prs.slide_height  # in EMUs

    tri_w = Inches(1.3)
    tri_h = Inches(0.75)
    gap   = Inches(0.05)

    # Right side: solid dark burgundy triangles flush to right edge
    for i in range(7):
        top = Inches(0.25) + i * (tri_h + gap)
        if top + tri_h > H:
            break
        tri = slide.shapes.add_shape(TRIANGLE, W - tri_w, top, tri_w, tri_h)
        tri.fill.solid()
        tri.fill.fore_color.rgb = dark_red if i % 2 == 0 else mid_red
        tri.line.fill.background()

    # Left side: white outline triangles flush to left edge
    for i in range(7):
        top = Inches(0.0) + i * (tri_h + gap)
        if top + tri_h > H:
            break
        tri = slide.shapes.add_shape(TRIANGLE, Inches(0), top, tri_w, tri_h)
        tri.fill.background()
        try:
            tri.line.color.rgb = white
            tri.line.width = Pt(1.5)
        except Exception:
            pass


def _add_triangles_closing(slide, prs):
    """Add white outline triangles flanking left and right edges for the closing slide."""
    white = RGBColor(255, 255, 255)
    W = prs.slide_width
    H = prs.slide_height

    tri_w = Inches(1.3)
    tri_h = Inches(0.75)
    gap = Inches(0.05)

    # Left side: white outline triangles flush to left edge
    for i in range(7):
        top = Inches(0.0) + i * (tri_h + gap)
        if top + tri_h > H:
            break
        tri = slide.shapes.add_shape(TRIANGLE, Inches(0), top, tri_w, tri_h)
        tri.fill.background()
        try:
            tri.line.color.rgb = white
            tri.line.width = Pt(1.5)
        except Exception:
            pass

    # Right side: white outline triangles flush to right edge
    for i in range(7):
        top = Inches(0.0) + i * (tri_h + gap)
        if top + tri_h > H:
            break
        tri = slide.shapes.add_shape(TRIANGLE, W - tri_w, top, tri_w, tri_h)
        tri.fill.background()
        try:
            tri.line.color.rgb = white
            tri.line.width = Pt(1.5)
        except Exception:
            pass


def _add_triangles_standard(slide, prs, has_right_image=False):
    """Add subtle gray triangle decorations for standard white slides - all within bounds."""
    light_gray = RGBColor(210, 210, 210)
    dark_gray  = RGBColor(70, 70, 70)
    W = prs.slide_width
    H = prs.slide_height

    tri_w = Inches(0.72)
    tri_h = Inches(0.60)

    # Right side: 3 light gray outline triangles flush to right edge - only if no image overlaps
    if not has_right_image:
        for i in range(3):
            top = Inches(2.4) + i * Inches(0.80)
            tri = slide.shapes.add_shape(TRIANGLE, W - tri_w, top, tri_w, tri_h)
            tri.fill.background()
            try:
                tri.line.color.rgb = light_gray
                tri.line.width = Pt(1.0)
            except Exception:
                pass

    # Bottom-left: 1 solid dark gray triangle, anchored at left edge
    bl_size = Inches(1.2)
    tri_bl = slide.shapes.add_shape(
        TRIANGLE,
        Inches(0),
        H - bl_size,
        bl_size, bl_size
    )
    tri_bl.fill.solid()
    tri_bl.fill.fore_color.rgb = dark_gray
    tri_bl.line.fill.background()
    set_shape_transparency(tri_bl, 25)



# ── Logo ──────────────────────────────────────────────────────────────

def _add_logo(slide, prs, left=None, top=None, dark_bg=False):
    """Add company logo. Renders transparent PNG directly on the background without any solid card/box."""
    logo_path = branding.LOGO_PATH

    if not logo_path.exists():
        assets_dir = logo_path.parent
        for fb_name in ["logo.png.png", "logo.png", "logo.webp"]:
            fb_path = assets_dir / fb_name
            if fb_path.exists():
                logo_path = fb_path
                break

    if not logo_path.exists():
        return

    try:
        if dark_bg:
            logo_h = Inches(1.35)
            logo_w = Inches(1.7)   # explicit width cap to stay in bounds
            if left is None:
                left = prs.slide_width - Inches(2.1)  # right-aligned within slide
            if top is None:
                top = Inches(0.25)
            # Render logo directly on the background
            slide.shapes.add_picture(str(logo_path), left, top, width=logo_w, height=logo_h)
        else:
            # Light background: render directly on top-left without any card/box
            logo_h = Inches(0.95)
            logo_w = Inches(1.2)
            if left is None:
                left = Inches(0.25)
            if top is None:
                top = Inches(0.15)
            slide.shapes.add_picture(str(logo_path), left, top, width=logo_w, height=logo_h)
    except Exception as e:
        logger.warning("Could not add logo: %s", e)



# ── Footer ───────────────────────────────────────────────────────────

def _add_footer(slide, prs, slide_num: int = None):
    """Three-part footer: page number left | company name centre | date right."""
    footer_top = prs.slide_height - Inches(0.48)
    footer_h = Inches(0.38)
    caption = branding.FONT_SIZE_CAPTION
    muted = branding.TEXT_COLOR_MUTED

    # Thin separator line
    sep = slide.shapes.add_shape(RECTANGLE, Inches(0.4), footer_top - Inches(0.06),
                                  prs.slide_width - Inches(0.8), Inches(0.01))
    sep.fill.solid()
    sep.fill.fore_color.rgb = RGBColor(210, 210, 210)
    sep.line.fill.background()

    # Left – page number
    if slide_num is not None:
        ln = slide.shapes.add_textbox(Inches(0.4), footer_top, Inches(1.0), footer_h)
        p = ln.text_frame.paragraphs[0]
        p.text = str(slide_num)
        p.font.size = caption
        p.font.color.rgb = muted
        p.alignment = PP_ALIGN.LEFT

    # Centre – company | report name
    cn = slide.shapes.add_textbox(Inches(3.0), footer_top, Inches(7.33), footer_h)
    p = cn.text_frame.paragraphs[0]
    p.text = f"{branding.FOOTER_TEXT_AR}  -  {branding.COMPANY_NAME_AR}"
    p.font.size = caption
    p.font.color.rgb = muted
    _mark_rtl(p)
    p.alignment = PP_ALIGN.CENTER

    # Right – date
    dr = slide.shapes.add_textbox(prs.slide_width - Inches(2.2), footer_top, Inches(1.8), footer_h)
    p = dr.text_frame.paragraphs[0]
    p.text = datetime.now().strftime("%d %B %Y")
    p.font.size = caption
    p.font.color.rgb = muted
    p.font.name = branding.FONT_ENGLISH
    p.alignment = PP_ALIGN.RIGHT


# ── Text Helpers ──────────────────────────────────────────────────────

def _add_title_shape(slide, prs, text: str, top: float = 0.45,
                     font_size=None, color=None, align=PP_ALIGN.RIGHT,
                     left: float = 0.5, width: float = None):
    """Add a styled title text box."""
    w = width if width is not None else (prs.slide_width / 914400 - left - 0.4)
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(1.35))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size or branding.FONT_SIZE_TITLE
    p.font.color.rgb = color or branding.PRIMARY_COLOR
    p.font.bold = True
    p.font.name = branding.FONT_HEADING
    _mark_rtl(p)
    p.alignment = align
    return txBox


def _add_body_text(slide, prs, text: str, top: float = 1.8,
                   left: float = 0.5, width: float = None):
    """Add body text with auto-scaling to prevent footer collisions."""
    w = width if width is not None else (prs.slide_width / 914400 - left - 0.5)
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(4.3))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    
    total_chars = len(text)
    
    # Calculate effective characters based on textbox width to scale font in narrow columns
    effective_w = w
    scale_factor = max(1.0, 8.0 / effective_w) if effective_w > 0 else 1.0
    scaled_chars = total_chars * scale_factor
    
    # More readable dynamic auto-scaling for body text
    if scaled_chars > 1600:
        font_size, spacing = Pt(9.5), Pt(1.0)
    elif scaled_chars > 1100:
        font_size, spacing = Pt(10.0), Pt(1.5)
    elif scaled_chars > 700:
        font_size, spacing = Pt(10.5), Pt(2.0)
    elif scaled_chars > 400:
        font_size, spacing = Pt(11.0), Pt(2.5)
    elif scaled_chars > 250:
        font_size, spacing = Pt(12.0), Pt(3.0)
    else:
        font_size, spacing = branding.FONT_SIZE_BODY, Pt(5)

    # Split paragraphs by newline to preserve paragraph separation and avoid huge line-breaks block
    paragraphs = [p.strip() for p in text.split("\n") if p.strip()]
    if not paragraphs:
        paragraphs = [""]
        
    for i, p_text in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = p_text
        p.font.size = font_size
        p.font.color.rgb = branding.TEXT_COLOR_DARK
        p.font.name = branding.FONT_BODY
        _mark_rtl(p)
        p.alignment = PP_ALIGN.RIGHT
        p.space_after = spacing

    return txBox


def _add_bullets(slide, prs, bullets: list, top: float = 1.8,
                 left: float = 0.5, width: float = None):
    """Add bulleted list with auto-scaling."""
    w = width if width is not None else (prs.slide_width / 914400 - left - 0.5)
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(4.3))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)

    total_chars = sum(len(b) for b in bullets)
    num = len(bullets)
    
    # Calculate effective characters based on textbox width to scale font in narrow columns
    effective_w = w
    scale_factor = max(1.0, 8.0 / effective_w) if effective_w > 0 else 1.0
    scaled_chars = total_chars * scale_factor
    
    # More readable dynamic auto-scaling for bullets
    if num > 9 or scaled_chars > 1600:
        font_size, spacing = Pt(9.5), Pt(1.0)
    elif num > 7 or scaled_chars > 1100:
        font_size, spacing = Pt(10.0), Pt(1.5)
    elif num > 5 or scaled_chars > 700:
        font_size, spacing = Pt(10.5), Pt(2.0)
    elif num > 4 or scaled_chars > 400:
        font_size, spacing = Pt(11.0), Pt(2.5)
    elif num > 3 or scaled_chars > 250:
        font_size, spacing = Pt(12.0), Pt(3.0)
    else:
        font_size, spacing = branding.FONT_SIZE_BODY, Pt(5)

    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"◄  {bullet}"
        p.font.size = font_size
        p.font.color.rgb = branding.TEXT_COLOR_DARK
        p.font.name = branding.FONT_BODY
        _mark_rtl(p)
        p.alignment = PP_ALIGN.RIGHT
        p.space_after = spacing

    return txBox


# ── Builders ──────────────────────────────────────────────────────────

def _build_cover(slide, prs, content: dict, image: Optional[bytes]):
    """Cover slide – Premium dark glassmorphic design.
    Uses high-resolution building background and a full-bleed semi-transparent dark mask
    allowing the background structure to remain visible underneath, with crisp white typography.
    """
    assets_dir = PROJECT_ROOT / "assets"
    bg_0_path = assets_dir / "bg_0.png"
    clean_bg_path = assets_dir / "clean_building.png"

    # Select best background image
    bg_path = None
    if clean_bg_path.exists():
        bg_path = clean_bg_path
    elif bg_0_path.exists():
        bg_path = bg_0_path

    # Prioritize dynamic custom generated image, fallback to static bg_path
    if image:
        _add_image(slide, prs, image, left=0, top=0)
    elif bg_path:
        _add_image_path(slide, prs, bg_path, left=0, top=0)

    if image or bg_path:
        # Layer 1: Base - Full-bleed building sunset photo (drawn above)

        # Layer 2: Full-bleed premium right-to-left gradient mask
        # Fades from rich brand orange #CC460A (90% opacity) on the right (text area) to 0% opacity on the left!
        grad_path = assets_dir / "cover_gradient.png"
        if grad_path.exists():
            _add_image_path(slide, prs, grad_path, left=0, top=0)
        else:
            full_mask = slide.shapes.add_shape(
                RECTANGLE,
                0,
                0,
                prs.slide_width,
                prs.slide_height
            )
            full_mask.fill.solid()
            full_mask.fill.fore_color.rgb = RGBColor(204, 70, 10) # Brand #CC460A
            set_shape_transparency(full_mask, 50) # 50% transparent (50% opacity)
            full_mask.line.fill.background()

        # Layer 3: Dynamic Brand Logo Card - top-right corner, clean and floating
        _add_logo(slide, prs, dark_bg=True)

        # Layer 4: Dynamic Title & Subtitle - right-aligned with 1.0 inch safe margin
        title_text = content.get("title", "دراسة الجدوى الأولية")
        subtitle_text = content.get("subtitle", "")

        # Title - above the divider
        txBox = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.33), Inches(2.0))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0)
        tf.margin_bottom = Inches(0)
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(42)
        p.font.color.rgb = RGBColor(255, 255, 255) # High-contrast pure white
        p.font.bold = True
        p.font.name = branding.FONT_HEADING
        _mark_rtl(p)
        p.alignment = PP_ALIGN.RIGHT

        # Thick Burgundy Divider Line (spanning the full width of the text block for structure)
        divider = slide.shapes.add_shape(
            RECTANGLE,
            Inches(1.0),
            Inches(4.4),
            Inches(11.33),
            Inches(0.06)
        )
        divider.fill.solid()
        divider.fill.fore_color.rgb = branding.PRIMARY_COLOR # Burgundy #670D0C
        divider.line.fill.background()

        # Subtitle - below the divider
        if subtitle_text:
            sub = slide.shapes.add_textbox(Inches(1.0), Inches(4.9), Inches(11.33), Inches(1.2))
            stf = sub.text_frame
            stf.word_wrap = True
            stf.margin_top = Inches(0)
            stf.margin_bottom = Inches(0)
            sp = stf.paragraphs[0]
            sp.text = subtitle_text
            sp.font.size = Pt(18)
            sp.font.color.rgb = RGBColor(220, 220, 220) # Premium light grey
            sp.font.name = branding.FONT_HEADING
            _mark_rtl(sp)
            sp.alignment = PP_ALIGN.RIGHT

        # Date - in the bottom-left corner
        db = slide.shapes.add_textbox(Inches(1.0), Inches(6.8), Inches(3.0), Inches(0.4))
        dtf = db.text_frame
        dtf.word_wrap = True
        dtf.margin_top = Inches(0)
        dtf.margin_bottom = Inches(0)
        dp = dtf.paragraphs[0]
        dp.text = datetime.now().strftime("%d %B %Y")
        dp.font.size = Pt(12)
        dp.font.color.rgb = RGBColor(255, 255, 255) # Pure white
        dp.font.name = branding.FONT_ENGLISH
        dp.alignment = PP_ALIGN.LEFT

    else:
        # Fallback dynamic design
        _set_slide_bg(slide, RGBColor(15, 15, 18)) # Dark base

        if image:
            _add_image(slide, prs, image, left=0, top=0,
                       width=prs.slide_width / 914400,
                       height=prs.slide_height / 914400)

        grad_path = assets_dir / "cover_gradient.png"
        if image and grad_path.exists():
            _add_image_path(slide, prs, grad_path, left=0, top=0)
        else:
            overlay = slide.shapes.add_shape(RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
            overlay.fill.solid()
            overlay.fill.fore_color.rgb = RGBColor(204, 70, 10) # #CC460A
            overlay.line.fill.background()
            set_shape_transparency(overlay, 50) # 50% opacity

        _add_triangles_cover(slide, prs)

        _add_logo(slide, prs, left=prs.slide_width - Inches(2.4), top=Inches(0.25), dark_bg=True)

        title_text = content.get("title", "دراسة الجدوى الأولية")
        subtitle_text = content.get("subtitle", "")

        txBox = slide.shapes.add_textbox(Inches(1.2), Inches(2.5), Inches(10.9), Inches(2.8))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0)
        tf.margin_bottom = Inches(0)
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(54)
        p.font.color.rgb = branding.TEXT_COLOR_LIGHT
        p.font.bold = True
        p.font.name = branding.FONT_HEADING
        _mark_rtl(p)
        p.alignment = PP_ALIGN.RIGHT

        divider = slide.shapes.add_shape(RECTANGLE, Inches(3.0), Inches(5.1),
                                          Inches(9.9), Inches(0.04))
        divider.fill.solid()
        divider.fill.fore_color.rgb = branding.SECONDARY_COLOR
        divider.line.fill.background()

        if subtitle_text:
            sub = slide.shapes.add_textbox(Inches(1.2), Inches(5.25), Inches(10.9), Inches(1.2))
            stf = sub.text_frame
            stf.word_wrap = True
            sp = stf.paragraphs[0]
            sp.text = subtitle_text
            sp.font.size = Pt(18)
            sp.font.color.rgb = branding.SECONDARY_COLOR
            sp.font.name = branding.FONT_HEADING
            _mark_rtl(sp)
            sp.alignment = PP_ALIGN.RIGHT

        db = slide.shapes.add_textbox(Inches(0.5), prs.slide_height - Inches(0.75),
                                       Inches(3.0), Inches(0.5))
        dp = db.text_frame.paragraphs[0]
        dp.text = datetime.now().strftime("%d %B %Y")
        dp.font.size = Pt(13)
        dp.font.color.rgb = branding.TEXT_COLOR_LIGHT
        dp.font.name = branding.FONT_ENGLISH
        dp.alignment = PP_ALIGN.LEFT


def _build_section_header(slide, prs, content: dict, image: Optional[bytes]):
    """Section header: full-bleed image + dark overlay + triangles + centered title."""
    assets_dir = Path(branding.LOGO_PATH).parent
    bg_4_path = assets_dir / "bg_4.png"
    clean_bg_path = assets_dir / "clean_building.png"
    section_overlay_path = assets_dir / "section_overlay_transparent.png"

    has_base_image = (image is not None) or clean_bg_path.exists() or bg_4_path.exists()

    if has_base_image and section_overlay_path.exists():
        # Layer 1: Base - Full-bleed custom image or clean building sunset photo
        if image:
            _add_image(slide, prs, image, left=0, top=0)
        elif clean_bg_path.exists():
            _add_image_path(slide, prs, clean_bg_path, left=0, top=0)
        elif bg_4_path.exists():
            _add_image_path(slide, prs, bg_4_path, left=0, top=0)

        # Layer 2: Overlay - Full-bleed transparent section overlay
        _add_image_path(slide, prs, section_overlay_path, left=0, top=0)

        title_text = content.get("title", "القسم")

        # Layer 3: Text - positioned exactly inside the beige card on the right
        # Section Label - above section title
        label = slide.shapes.add_textbox(Inches(6.15), Inches(3.1), Inches(6.2), Inches(0.4))
        ltf = label.text_frame
        ltf.word_wrap = True
        ltf.margin_top = Inches(0)
        ltf.margin_bottom = Inches(0)
        lp = ltf.paragraphs[0]
        lp.text = "◈  قسم"
        lp.font.size = Pt(14)
        lp.font.color.rgb = branding.SECONDARY_COLOR
        lp.font.bold = True
        lp.font.name = branding.FONT_HEADING
        _mark_rtl(lp)
        lp.alignment = PP_ALIGN.RIGHT

        # Section Title - large white title inside the orange gradient
        txBox = slide.shapes.add_textbox(Inches(6.15), Inches(3.55), Inches(6.2), Inches(0.6))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0)
        tf.margin_bottom = Inches(0)
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(36)
        p.font.color.rgb = RGBColor(255, 255, 255) # Pure White for high legibility
        p.font.bold = True
        p.font.name = branding.FONT_HEADING
        _mark_rtl(p)
        p.alignment = PP_ALIGN.RIGHT

    elif image or bg_4_path.exists():
        if image:
            _add_image(slide, prs, image, left=0, top=0)
        else:
            _add_image_path(slide, prs, bg_4_path, left=0, top=0)

        title_text = content.get("title", "القسم")

        # Section label
        label = slide.shapes.add_textbox(Inches(6.15), Inches(3.1), Inches(6.2), Inches(0.4))
        ltf = label.text_frame
        ltf.word_wrap = True
        ltf.margin_top = Inches(0)
        ltf.margin_bottom = Inches(0)
        lp = ltf.paragraphs[0]
        lp.text = "◈  قسم"
        lp.font.size = Pt(14)
        lp.font.color.rgb = branding.SECONDARY_COLOR
        lp.font.bold = True
        lp.font.name = branding.FONT_HEADING
        _mark_rtl(lp)
        lp.alignment = PP_ALIGN.RIGHT

        # Title inside the orange gradient (pure white for high legibility)
        txBox = slide.shapes.add_textbox(Inches(6.15), Inches(3.55), Inches(6.2), Inches(0.6))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0)
        tf.margin_bottom = Inches(0)
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(36)
        p.font.color.rgb = RGBColor(255, 255, 255) # Pure White
        p.font.bold = True
        p.font.name = branding.FONT_HEADING
        _mark_rtl(p)
        p.alignment = PP_ALIGN.RIGHT
    else:
        # Fallback dynamic design
        _set_slide_bg(slide, branding.PRIMARY_COLOR)

        if image:
            _add_image(slide, prs, image, left=0, top=0,
                       width=prs.slide_width / 914400,
                       height=prs.slide_height / 914400)

        overlay = slide.shapes.add_shape(RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = branding.PRIMARY_COLOR
        overlay.line.fill.background()
        set_shape_transparency(overlay, 38)

        _add_triangles_cover(slide, prs)

        _add_logo(slide, prs, left=Inches(0.35), top=Inches(0.25), dark_bg=True)

        title_text = content.get("title", "القسم")
        label = slide.shapes.add_textbox(Inches(1.2), Inches(2.8), Inches(11.0), Inches(0.55))
        lp = label.text_frame.paragraphs[0]
        lp.text = "◈  قسم"
        lp.font.size = Pt(14)
        lp.font.color.rgb = branding.SECONDARY_COLOR
        lp.font.bold = True
        lp.font.name = branding.FONT_HEADING
        _mark_rtl(lp)
        lp.alignment = PP_ALIGN.RIGHT

        txBox = slide.shapes.add_textbox(Inches(1.2), Inches(3.25), Inches(11.0), Inches(2.2))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0)
        tf.margin_bottom = Inches(0)
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(52)
        p.font.color.rgb = branding.TEXT_COLOR_LIGHT
        p.font.bold = True
        p.font.name = branding.FONT_HEADING
        _mark_rtl(p)
        p.alignment = PP_ALIGN.RIGHT

        acc = slide.shapes.add_shape(RECTANGLE, Inches(3.0), Inches(5.4),
                                      Inches(9.9), Inches(0.04))
        acc.fill.solid()
        acc.fill.fore_color.rgb = branding.SECONDARY_COLOR
        acc.line.fill.background()


def _build_standard_slide(slide, prs, content: dict, image: Optional[bytes]):
    """Standard content slide: white bg + burgundy title + triangle decorations."""
    _set_slide_bg(slide, branding.TEXT_COLOR_LIGHT)

    slide_idx = content.get("slide_index", 0)
    has_image = bool(image)
    has_right_image = has_image and (slide_idx % 2 != 0)

    # Triangle decorations (reference style)
    _add_triangles_standard(slide, prs, has_right_image=has_right_image)

    # Logo – top left without card
    _add_logo(slide, prs, dark_bg=False)

    # Title – large burgundy, right-aligned
    title = content.get("title", "")
    _add_title_shape(slide, prs, title, top=0.45, font_size=Pt(38),
                     color=branding.PRIMARY_COLOR, left=1.9,
                     width=prs.slide_width / 914400 - 2.4)

    # Thin gold underline under title
    und = slide.shapes.add_shape(RECTANGLE, Inches(1.9), Inches(1.5),
                                  Inches(11.0), Inches(0.03))
    und.fill.solid()
    und.fill.fore_color.rgb = branding.SECONDARY_COLOR
    und.line.fill.background()

    # Content layout: alternating layout
    has_bullets = "bullets" in content or "highlights" in content or "key_features" in content

    # Slide dimensions: 13.33" wide, 7.5" tall
    # Image zone: left or right → width=5.5", height=5.0"
    IMG_TOP   = 1.55
    IMG_W     = 5.5
    IMG_H     = 5.0

    if image:
        if slide_idx % 2 == 0:
            # Even slide index: Image on LEFT, Text on RIGHT
            IMG_LEFT  = 0.4
            TEXT_LEFT = 6.3
            TEXT_W    = 6.6
        else:
            # Odd slide index: Image on RIGHT, Text on LEFT
            IMG_LEFT  = 7.4
            TEXT_LEFT = 0.4
            TEXT_W    = 6.6
            
        _add_image(slide, prs, image, left=IMG_LEFT, top=IMG_TOP, width=IMG_W, height=IMG_H)
        t_left, t_w = TEXT_LEFT, TEXT_W
    else:
        # No image → text spans full width (minus triangle zone on right)
        t_left, t_w = 0.4, 11.5

    if has_bullets:
        bullets = content.get("bullets", content.get("highlights", content.get("key_features", [])))
        if isinstance(bullets, list) and bullets:
            _add_bullets(slide, prs, bullets, top=1.65, left=t_left, width=t_w)
        else:
            desc = content.get("description", content.get("concept", content.get("summary", "")))
            if desc:
                _add_body_text(slide, prs, desc, top=1.65, left=t_left, width=t_w)
    else:
        desc = content.get("description", content.get("concept", content.get("message", "")))
        if desc:
            _add_body_text(slide, prs, desc, top=1.65, left=t_left, width=t_w)

    _add_footer(slide, prs, slide_num=slide_idx + 1)



def _build_two_column_slide(slide, prs, content: dict, image: Optional[bytes]):
    _build_standard_slide(slide, prs, content, image)


def _build_timeline_slide(slide, prs, content: dict, image: Optional[bytes]):
    """Timeline slide with auto-scaling phase entries and RTL markers."""
    _set_slide_bg(slide, branding.TEXT_COLOR_LIGHT)
    _add_triangles_standard(slide, prs)
    _add_logo(slide, prs, dark_bg=False)

    title = content.get("title", "الجدول الزمني")
    _add_title_shape(slide, prs, title, top=0.45, font_size=Pt(38),
                     color=branding.PRIMARY_COLOR, left=1.9)

    # Gold underline
    und = slide.shapes.add_shape(RECTANGLE, Inches(1.9), Inches(1.5),
                                  Inches(11.0), Inches(0.03))
    und.fill.solid()
    und.fill.fore_color.rgb = branding.SECONDARY_COLOR
    und.line.fill.background()

    phases = content.get("phases", [])
    display = phases[:5]
    n = len(display)
    step = 0.80 if n >= 5 else (0.88 if n == 4 else 1.0)
    phase_fs = Pt(15) if n >= 5 else (Pt(17) if n == 4 else branding.FONT_SIZE_SUBHEADING)
    desc_fs = Pt(13) if n >= 5 else (Pt(14) if n == 4 else branding.FONT_SIZE_BODY)

    top = 1.8
    for phase in display:
        name = phase.get("name", "")
        duration = phase.get("duration", "")
        desc = phase.get("description", "")

        txBox = slide.shapes.add_textbox(Inches(1.5), Inches(top), Inches(10.5), Inches(step))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0)
        tf.margin_bottom = Inches(0)
        p = tf.paragraphs[0]
        p.text = f"◆  {name}  ({duration})"
        p.font.size = phase_fs
        p.font.color.rgb = branding.PRIMARY_COLOR
        p.font.bold = True
        p.font.name = branding.FONT_HEADING
        _mark_rtl(p)
        p.alignment = PP_ALIGN.RIGHT

        p2 = tf.add_paragraph()
        p2.text = f"       {desc}"
        p2.font.size = desc_fs
        p2.font.color.rgb = branding.TEXT_COLOR_DARK
        p2.font.name = branding.FONT_BODY
        _mark_rtl(p2)
        p2.alignment = PP_ALIGN.RIGHT
        p2.space_after = Pt(2)
        top += step

    _add_footer(slide, prs, slide_num=content.get("slide_index", 0) + 1)


def _build_closing_slide(slide, prs, content: dict, image: Optional[bytes]):
    """Closing slide matching reference 41.png exactly or adjusting dynamically if content is long.
    Solid dark burgundy background, centered logo, white outline triangles flanking left/right,
    and centered contact info at the bottom.
    """
    # 1. Set solid dark burgundy background panel (#670D0C)
    _set_slide_bg(slide, branding.PRIMARY_COLOR)

    # 2. Add white outline triangles flanking left and right edges
    _add_triangles_closing(slide, prs)

    message = content.get("message", content.get("title", "شكراً لكم لتواصلكم معنا"))
    is_long = len(message) > 150

    # 3. Add logo - scaled & positioned dynamically based on text length
    logo_path = branding.LOGO_PATH
    if not logo_path.exists():
        assets_dir = logo_path.parent
        for fb_name in ["logo.png.png", "logo.png", "logo.webp"]:
            fb_path = assets_dir / fb_name
            if fb_path.exists():
                logo_path = fb_path
                break

    if logo_path.exists():
        try:
            if is_long:
                # Long text: logo is smaller and pushed to the top
                logo_w = Inches(1.6)
                logo_h = Inches(1.2)
                logo_left = (prs.slide_width - logo_w) / 2
                logo_top = Inches(0.3)
            else:
                # Short text (standard): logo is larger and centered vertically
                logo_w = Inches(2.4)
                logo_h = Inches(1.9)
                logo_left = (prs.slide_width - logo_w) / 2
                logo_top = (prs.slide_height - logo_h) / 2 - Inches(0.5)

            # Draw PNG directly on the closing slide background without any solid card
            slide.shapes.add_picture(str(logo_path), logo_left, logo_top, width=logo_w, height=logo_h)
        except Exception as e:
            logger.warning("Could not add closing slide logo: %s", e)

    # 4. Message block and contact details
    if is_long:
        # Dynamic detailed summary layout
        txBox = slide.shapes.add_textbox(Inches(1.2), Inches(1.6), prs.slide_width - Inches(2.4), Inches(4.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0)
        tf.margin_bottom = Inches(0)

        # Dynamic scaling for long closing message
        total_chars = len(message)
        if total_chars > 800:
            font_size = Pt(9.5)
            spacing = Pt(1)
        elif total_chars > 500:
            font_size = Pt(10.5)
            spacing = Pt(2)
        elif total_chars > 300:
            font_size = Pt(11.5)
            spacing = Pt(3)
        else:
            font_size = Pt(13.0)
            spacing = Pt(4)

        # Split message by paragraph to preserve paragraph separation
        paragraphs = [p.strip() for p in message.split("\n") if p.strip()]
        for i, p_text in enumerate(paragraphs):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = p_text
            p.font.size = font_size
            p.font.bold = (i == 0 and total_chars < 500) # Bold first paragraph if short-ish
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.name = branding.FONT_HEADING if i == 0 else branding.FONT_BODY
            _mark_rtl(p)
            p.alignment = PP_ALIGN.CENTER
            p.space_after = spacing

        # Draw contact details in a separate absolute-positioned textbox at the very bottom
        contact_top = prs.slide_height - Inches(0.8)
        contact_tx = slide.shapes.add_textbox(Inches(1.2), contact_top, prs.slide_width - Inches(2.4), Inches(0.5))
        ctf = contact_tx.text_frame
        ctf.word_wrap = True
        ctf.margin_top = Inches(0)
        ctf.margin_bottom = Inches(0)
        cp = ctf.paragraphs[0]
        cp.text = f"{branding.COMPANY_NAME_AR}  |  info@manafe.com.sa  |  www.manafe.com.sa"
        cp.font.size = Pt(12)
        cp.font.color.rgb = RGBColor(220, 220, 220)
        cp.font.name = branding.FONT_BODY
        _mark_rtl(cp)
        cp.alignment = PP_ALIGN.CENTER

    else:
        # Standard thank you layout (original design)
        contact_top = prs.slide_height - Inches(2.0)
        contact_h = Inches(1.5)
        contact_w = prs.slide_width - Inches(2.0)
        contact_left = Inches(1.0)

        txBox = slide.shapes.add_textbox(contact_left, contact_top, contact_w, contact_h)
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0)
        tf.margin_bottom = Inches(0)

        p1 = tf.paragraphs[0]
        p1.text = message
        p1.font.size = Pt(20)
        p1.font.bold = True
        p1.font.color.rgb = RGBColor(255, 255, 255)
        p1.font.name = branding.FONT_HEADING
        _mark_rtl(p1)
        p1.alignment = PP_ALIGN.CENTER
        p1.space_after = Pt(10)

        p2 = tf.add_paragraph()
        p2.text = f"{branding.COMPANY_NAME_AR}  |  info@manafe.com.sa  |  www.manafe.com.sa"
        p2.font.size = Pt(14)
        p2.font.color.rgb = RGBColor(220, 220, 220)
        p2.font.name = branding.FONT_BODY
        _mark_rtl(p2)
        p2.alignment = PP_ALIGN.CENTER


def _build_swot_slide(slide, prs, content: dict, image: Optional[bytes]):
    _set_slide_bg(slide, branding.TEXT_COLOR_LIGHT)
    _add_triangles_standard(slide, prs)
    _add_logo(slide, prs, dark_bg=False)
    title = content.get("title", "تحليل SWOT")
    _add_title_shape(slide, prs, title, top=0.45, font_size=Pt(38),
                     color=branding.PRIMARY_COLOR, left=1.9)

    und = slide.shapes.add_shape(RECTANGLE, Inches(1.9), Inches(1.5),
                                  Inches(11.0), Inches(0.03))
    und.fill.solid()
    und.fill.fore_color.rgb = branding.SECONDARY_COLOR
    und.line.fill.background()

    silver = RGBColor(167, 169, 172)
    sections = [
        ("نقاط القوة",   content.get("strengths", []),     branding.PRIMARY_COLOR,   9.2),
        ("نقاط الضعف",   content.get("weaknesses", []),    branding.SECONDARY_COLOR, 6.4),
        ("الفرص",        content.get("opportunities", []), silver,                    3.6),
        ("التحديات",     content.get("threats", []),       branding.PRIMARY_COLOR,   0.8),
    ]

    for label, points, color, left in sections:
        header = slide.shapes.add_shape(RECTANGLE, Inches(left), Inches(1.7), Inches(2.7), Inches(0.75))
        header.fill.solid()
        header.fill.fore_color.rgb = color
        header.line.fill.background()
        htf = header.text_frame
        hp = htf.paragraphs[0]
        hp.text = label
        hp.font.size = branding.FONT_SIZE_HEADING
        hp.font.color.rgb = branding.TEXT_COLOR_LIGHT
        hp.font.bold = True
        hp.font.name = branding.FONT_HEADING
        hp.alignment = PP_ALIGN.CENTER

        _add_bullets(slide, prs, points, top=2.6, left=left, width=2.7)

    _add_footer(slide, prs, slide_num=content.get("slide_index", 0) + 1)


def _build_map_slide(slide, prs, content: dict, image: Optional[bytes]):
    _set_slide_bg(slide, branding.TEXT_COLOR_LIGHT)
    _add_triangles_standard(slide, prs, has_right_image=bool(image))
    _add_logo(slide, prs, dark_bg=False)

    title = content.get("title", "تحليل الموقع")
    _add_title_shape(slide, prs, title, top=0.45, font_size=Pt(38),
                     color=branding.PRIMARY_COLOR, left=1.9)

    und = slide.shapes.add_shape(RECTANGLE, Inches(1.9), Inches(1.5),
                                  Inches(11.0), Inches(0.03))
    und.fill.solid()
    und.fill.fore_color.rgb = branding.SECONDARY_COLOR
    und.line.fill.background()

    if image:
        _add_image(slide, prs, image, left=4.2, top=1.65, width=8.7, height=5.5)

    desc = content.get("description", content.get("desc", ""))
    bullets = content.get("bullets", [])
    sidebar_color = branding.PRIMARY_COLOR

    # Left sidebar with dark panel
    panel = slide.shapes.add_shape(RECTANGLE, 0, 0, Inches(4.0), prs.slide_height)
    panel.fill.solid()
    panel.fill.fore_color.rgb = sidebar_color
    panel.line.fill.background()
    set_shape_transparency(panel, 8)

    if desc or bullets:
        # Combine desc and bullets into a single textbox to prevent overlaps
        txBox = slide.shapes.add_textbox(Inches(0.2), Inches(1.8), Inches(3.6), Inches(5.0))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0)
        tf.margin_bottom = Inches(0)

        # Calculate total characters and count to determine dynamic scaling
        total_chars = len(desc) + sum(len(b) for b in bullets)
        num_items = len(bullets) + (1 if desc else 0)
        
        # Apply narrow column scaling factor (width = 3.6, so 8.0/3.6 = 2.22)
        scale_factor = 2.22
        scaled_chars = total_chars * scale_factor
        
        # More readable dynamic auto-scaling for sidebar content
        if num_items > 7 or scaled_chars > 1400:
            font_size, spacing = Pt(9.0), Pt(0.5)
        elif num_items > 5 or scaled_chars > 1000:
            font_size, spacing = Pt(9.5), Pt(1.0)
        elif num_items > 3 or scaled_chars > 600:
            font_size, spacing = Pt(10.0), Pt(1.5)
        elif scaled_chars > 300:
            font_size, spacing = Pt(10.5), Pt(2.0)
        else:
            font_size, spacing = Pt(11.5), Pt(3.0)

        is_first = True

        if desc:
            paragraphs = [p.strip() for p in desc.split("\n") if p.strip()]
            for p_text in paragraphs:
                p = tf.paragraphs[0] if is_first else tf.add_paragraph()
                is_first = False
                p.text = p_text
                p.font.size = font_size
                p.font.color.rgb = branding.TEXT_COLOR_LIGHT
                p.font.name = branding.FONT_BODY
                _mark_rtl(p)
                p.alignment = PP_ALIGN.RIGHT
                p.space_after = spacing

        if bullets:
            if not is_first:
                # Add a bit of space before bullets if there was description
                tf.paragraphs[-1].space_after = spacing + Pt(4)
                
            for bullet in bullets:
                p = tf.paragraphs[0] if is_first else tf.add_paragraph()
                is_first = False
                p.text = f"◄  {bullet}"
                p.font.size = font_size
                p.font.color.rgb = branding.TEXT_COLOR_LIGHT
                p.font.name = branding.FONT_BODY
                _mark_rtl(p)
                p.alignment = PP_ALIGN.RIGHT
                p.space_after = spacing

    _add_footer(slide, prs, slide_num=content.get("slide_index", 0) + 1)


def style_cell(cell, bg_color, text, font_size, text_color, bold=False, align=PP_ALIGN.RIGHT):
    """Style a PowerPoint table cell with a specific background, padding, word wrap, and font settings."""
    cell.fill.solid()
    cell.fill.fore_color.rgb = bg_color
    cell.text_frame.margin_left = Inches(0.15)
    cell.text_frame.margin_right = Inches(0.15)
    cell.text_frame.margin_top = Inches(0.08)
    cell.text_frame.margin_bottom = Inches(0.08)
    cell.text_frame.word_wrap = True
    
    p = cell.text_frame.paragraphs[0]
    p.text = text
    p.font.name = branding.FONT_BODY
    p.font.size = font_size
    p.font.color.rgb = text_color
    p.font.bold = bold
    _mark_rtl(p)
    p.alignment = align


def _build_chart_slide(slide, prs, content: dict, image: Optional[bytes]):
    """Financial slide: Premium Native PowerPoint Table layout for financial allocations/data."""
    _set_slide_bg(slide, branding.TEXT_COLOR_LIGHT)
    _add_triangles_standard(slide, prs)
    _add_logo(slide, prs, dark_bg=False)
    title = content.get("title", "البيانات المالية")
    _add_title_shape(slide, prs, title, top=0.45, font_size=Pt(38),
                     color=branding.PRIMARY_COLOR, left=1.9)

    und = slide.shapes.add_shape(RECTANGLE, Inches(1.9), Inches(1.5),
                                  Inches(11.0), Inches(0.03))
    und.fill.solid()
    und.fill.fore_color.rgb = branding.SECONDARY_COLOR
    und.line.fill.background()

    items = content.get("items") or content.get("data") or []
    
    if items:
        # Limit rows to 8 to avoid spilling over slide margins/footer
        items_to_show = items[:8]
        rows = len(items_to_show) + 1
        cols = 2
        
        # Position table perfectly below title & line
        t_left = Inches(1.0)
        t_top = Inches(1.8)
        t_width = Inches(11.33)
        t_height = Inches(min(4.8, 0.55 * rows))
        
        table_shape = slide.shapes.add_table(rows, cols, t_left, t_top, t_width, t_height)
        table = table_shape.table
        table.columns[0].width = Inches(3.5)   # Left column: values
        table.columns[1].width = Inches(7.83)  # Right column: descriptions
        
        # Header Row
        style_cell(table.cell(0, 0), branding.PRIMARY_COLOR, "القيمة / التقدير", Pt(14), RGBColor(255, 255, 255), bold=True, align=PP_ALIGN.CENTER)
        style_cell(table.cell(0, 1), branding.PRIMARY_COLOR, "البند / البيان المالي", Pt(14), RGBColor(255, 255, 255), bold=True, align=PP_ALIGN.RIGHT)
        
        # Data Rows
        for r_idx, item in enumerate(items_to_show, start=1):
            label = item.get("label", "")
            value = str(item.get("value", ""))
            
            # Alternate row background colors: light beige vs pure white
            bg_col = RGBColor(245, 240, 238) if r_idx % 2 == 1 else RGBColor(255, 255, 255)
            
            style_cell(table.cell(r_idx, 0), bg_col, value, Pt(12), branding.TEXT_COLOR_DARK, bold=True, align=PP_ALIGN.CENTER)
            style_cell(table.cell(r_idx, 1), bg_col, label, Pt(12), branding.TEXT_COLOR_DARK, bold=False, align=PP_ALIGN.RIGHT)
    else:
        # Fallback description text if no items are present
        desc = content.get("description", "")
        if desc:
            _add_body_text(slide, prs, desc, top=1.8, left=1.0, width=11.33)

    _add_footer(slide, prs, slide_num=content.get("slide_index", 0) + 1)


# ── Dispatch Table ────────────────────────────────────────────────────

SLIDE_BUILDERS = {
    "cover":          _build_cover,
    "section_header": _build_section_header,
    "content":        _build_standard_slide,
    "two_column":     _build_two_column_slide,
    "map":            _build_map_slide,
    "swot":           _build_swot_slide,
    "chart":          _build_chart_slide,
    "timeline":       _build_timeline_slide,
    "closing":        _build_closing_slide,
}


def build_presentation(
    slide_contents: list,
    images: dict,
    project_name: str = "proposal",
    theme: dict = None,
) -> Path:
    logger.info("Building presentation with %d slides...", len(slide_contents))

    if theme:
        if "bg_color" in theme:
            branding.BACKGROUND_COLOR = _get_color(theme["bg_color"])
        if "text_color" in theme:
            branding.TEXT_COLOR_DARK = _get_color(theme["text_color"])
        if "font" in theme:
            branding.FONT_HEADING = theme["font"]
            branding.FONT_BODY = theme["font"]

    template_path = settings.TEMPLATE_PATH
    prs = Presentation(str(template_path)) if template_path.exists() else Presentation()

    prs.slide_width  = Inches(branding.SLIDE_WIDTH_INCHES)
    prs.slide_height = Inches(branding.SLIDE_HEIGHT_INCHES)

    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]

    SLIDE_TYPE_MAP = {
        "cover":          "cover",
        "section_header": "section_header",
        "timeline":       "timeline",
        "swot":           "swot",
        "map":            "map",
        "two_column":     "two_column",
        "chart":          "chart",
        "standard":       "content",
        "closing":        "closing",
    }

    for slide_data in slide_contents:
        slide_type  = slide_data.get("slide_type", "standard")
        layout_type = SLIDE_TYPE_MAP.get(slide_type, "content")
        slide_idx   = slide_data.get("slide_index", 0)
        content     = slide_data
        image       = images.get(slide_idx)

        # Fallback to the gorgeous cover image for section divider slides to maintain brand consistency
        if layout_type == "section_header" and not image:
            image = images.get(0)

        slide = prs.slides.add_slide(blank_layout)
        builder = SLIDE_BUILDERS.get(layout_type, _build_standard_slide)
        builder(slide, prs, content, image)

    timestamp  = datetime.now().strftime("%Y%m%d_%H%M%S")
    safe_name  = "".join(c if c.isalnum() or c in "-_ " else "_" for c in project_name)
    output_path = settings.OUTPUT_DIR / f"{safe_name}_{timestamp}.pptx"
    prs.save(str(output_path))
    logger.info("Presentation saved to %s", output_path)
    return output_path
