import fitz  # PyMuPDF
import docx
from pptx import Presentation
import io
import logging
import base64
from typing import List, Any, Tuple

logger = logging.getLogger(__name__)

def extract_text_from_file(file_bytes: bytes, filename: str) -> str:
    """Extract text from supported file types."""
    ext = filename.lower().split('.')[-1]
    
    try:
        if ext == 'txt':
            return file_bytes.decode('utf-8')
            
        elif ext == 'pdf':
            doc = fitz.open(stream=file_bytes, filetype="pdf")
            text = []
            for page in doc:
                text.append(page.get_text())
            return "\n".join(text)
            
        elif ext in ['doc', 'docx']:
            doc = docx.Document(io.BytesIO(file_bytes))
            text = []
            for para in doc.paragraphs:
                text.append(para.text)
            return "\n".join(text)
            
        elif ext == 'pptx':
            prs = Presentation(io.BytesIO(file_bytes))
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return "\n".join(text)
            
        else:
            logger.warning(f"Unsupported file type for text extraction: {ext}")
            return ""
            
    except Exception as e:
        logger.error(f"Error extracting text from {filename}: {e}")
        return f"[Error extracting {filename}]"

def extract_base64_image(file_bytes: bytes) -> str:
    """Resize image to max 800px width/height and compress to JPEG to minimize base64 payload."""
    try:
        from PIL import Image
        img = Image.open(io.BytesIO(file_bytes))
        if img.mode in ("RGBA", "P"):
            img = img.convert("RGB")
        img.thumbnail((800, 800), Image.Resampling.LANCZOS)
        out_buf = io.BytesIO()
        img.save(out_buf, format="JPEG", quality=75)
        compressed = out_buf.getvalue()
        logger.info("Compressed image payload: from %d bytes to %d bytes", len(file_bytes), len(compressed))
        return base64.b64encode(compressed).decode('utf-8')
    except Exception as e:
        logger.warning("Failed to compress image, using original: %s", e)
        return base64.b64encode(file_bytes).decode('utf-8')

def process_uploaded_files(uploaded_files: List[Any]) -> Tuple[str, List[dict]]:
    """
    Process a list of Streamlit UploadedFile objects.
    Returns:
        tuple: (concatenated_text, list_of_image_dicts)
        where image_dicts are {"filename": str, "base64": str, "mime_type": str}
    """
    if not uploaded_files:
        return "", []
        
    combined_text = []
    images = []
    
    for file in uploaded_files:
        ext = file.name.lower().split('.')[-1]
        file_bytes = file.getvalue()
        
        if ext in ['png', 'jpg', 'jpeg']:
            images.append({
                "filename": file.name,
                "base64": extract_base64_image(file_bytes),
                "mime_type": "image/jpeg"
            })
        else:
            text = extract_text_from_file(file_bytes, file.name)
            if text.strip():
                combined_text.append(f"--- Document: {file.name} ---\n{text.strip()}\n")
            
    return "\n".join(combined_text), images
